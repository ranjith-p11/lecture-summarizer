import os

def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from a PDF, DOCX, or PPTX file.
    Returns the extracted text as a string.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return _extract_from_pdf(file_path)
    elif ext == ".docx":
        return _extract_from_docx(file_path)
    elif ext == ".pptx":
        return _extract_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported document format: {ext}")

def _extract_from_pdf(file_path: str) -> str:
    import PyPDF2
    text_content = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_content.append(page_text)
    return "\n".join(text_content).strip()

def _extract_from_docx(file_path: str) -> str:
    import docx
    doc = docx.Document(file_path)
    text_content = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(text_content).strip()

def _extract_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    text_content = []
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
    return "\n".join(text_content).strip()
